# file_handler.py
"""Gestionnaire pour l'extraction de contenu des fichiers"""

import PyPDF2
from docx import Document
import pandas as pd
from pptx import Presentation
import openpyxl
import base64
from PIL import Image
import re
import hashlib
import io
from config import IMAGE_EXTENSIONS, IMAGE_MIME_TYPES

class FileHandler:
    """Classe pour gérer l'extraction de contenu des fichiers"""
    
    @staticmethod
    def extract_text(file):
        """Extrait le texte des fichiers uploadés avec support complet Office + Images"""
        try:
            # Sauvegarder la position actuelle du fichier
            original_position = file.tell() if hasattr(file, 'tell') else 0
            
            file_type = file.type.lower() if hasattr(file, 'type') else ""
            file_name = file.name.lower() if hasattr(file, 'name') else ""
            
            # Images - retourner les métadonnées
            if FileHandler.is_image_file(file):
                result = FileHandler._extract_image_metadata(file)
            
            # PDF
            elif file_type == "application/pdf":
                result = FileHandler._extract_pdf_text(file)
            
            # Word Documents (.docx)
            elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                result = FileHandler._extract_docx_text(file)
            
            # Word Documents anciens (.doc)
            elif file_name.endswith('.doc'):
                result = FileHandler._extract_doc_text(file)
            
            # Excel Files (.xlsx)
            elif file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                result = FileHandler._extract_xlsx_text(file)
            
            # Excel anciens (.xls)
            elif file_name.endswith('.xls'):
                result = FileHandler._extract_xls_text(file)
            
            # PowerPoint (.pptx)
            elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                result = FileHandler._extract_pptx_text(file)
            
            # PowerPoint ancien (.ppt)
            elif file_name.endswith('.ppt'):
                result = "Format .ppt non supporté. Veuillez convertir en .pptx"
            
            # CSV
            elif file_type == "text/csv" or file_name.endswith('.csv'):
                result = FileHandler._extract_csv_text(file)
            
            # RTF
            elif file_name.endswith('.rtf'):
                result = FileHandler._extract_rtf_text(file)
            
            # Texte plain
            elif file_type == "text/plain":
                file.seek(0)
                content = file.read()
                if isinstance(content, bytes):
                    result = content.decode("utf-8", errors='ignore')
                else:
                    result = str(content)
            
            else:
                result = f"Type de fichier non supporté: {file_type}"
            
            # Restaurer la position du fichier
            try:
                file.seek(original_position)
            except:
                pass
                
            return result
                
        except Exception as e:
            return f"Erreur lors de l'extraction: {str(e)}"
    
    @staticmethod
    def _extract_image_metadata(file):
        """Extrait les métadonnées d'une image"""
        try:
            file.seek(0)
            image = Image.open(file)
            info = {
                "format": image.format,
                "mode": image.mode,
                "size": image.size,
                "filename": getattr(file, 'name', 'image')
            }
            
            file_size = len(file.getvalue()) if hasattr(file, 'getvalue') else getattr(file, 'size', 0)
            
            metadata_text = f"""Image: {info['filename']}
Format: {info['format']}
Dimensions: {info['size'][0]} x {info['size'][1]} pixels
Mode couleur: {info['mode']}
Taille fichier: {FileHandler.format_file_size(file_size)}
"""
            return metadata_text
        except Exception as e:
            return f"Erreur lecture image: {str(e)}"
    
    @staticmethod
    def _extract_pdf_text(file):
        """Extrait le texte d'un fichier PDF"""
        try:
            file.seek(0)
            reader = PyPDF2.PdfReader(file)
            text = ""
            for page_num, page in enumerate(reader.pages):
                try:
                    text += f"\n--- Page {page_num + 1} ---\n"
                    page_text = page.extract_text()
                    text += page_text or ""
                except Exception as e:
                    text += f"\n[Erreur lecture page {page_num + 1}: {str(e)}]\n"
            return text
        except Exception as e:
            return f"Erreur extraction PDF: {str(e)}"
    
    @staticmethod
    def _extract_docx_text(file):
        """Extrait le texte d'un fichier Word .docx"""
        try:
            file.seek(0)
            doc = Document(file)
            text = ""
            
            # Extraire les paragraphes
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
            
            # Extraire le texte des tableaux
            if doc.tables:
                for table in doc.tables:
                    text += "\n--- Tableau ---\n"
                    for row in table.rows:
                        row_text = " | ".join([cell.text.strip() for cell in row.cells])
                        if row_text.strip():
                            text += row_text + "\n"
            
            return text or "Document vide"
        except Exception as e:
            return f"Erreur extraction DOCX: {str(e)}"
    
    @staticmethod
    def _extract_doc_text(file):
        """Extraction basique pour les fichiers .doc (format binaire ancien)"""
        try:
            file.seek(0)
            content = file.read()
            if isinstance(content, bytes):
                text = content.decode('utf-8', errors='ignore')
            else:
                text = str(content)
            
            # Nettoyer le texte
            text = re.sub(r'[\x00-\x1f\x7f-\x9f]', ' ', text)
            words = [word for word in text.split() if len(word) > 1]
            return ' '.join(words[:1000]) if words else "Impossible d'extraire le texte du fichier .doc"
        except Exception as e:
            return f"Erreur extraction DOC: {str(e)}"
    
    @staticmethod
    def _extract_xlsx_text(file):
        """Extrait le texte d'un fichier Excel .xlsx"""
        try:
            file.seek(0)
            workbook = openpyxl.load_workbook(file, data_only=True, read_only=True)
            text = ""
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"\n--- Feuille: {sheet_name} ---\n"
                
                for row in sheet.iter_rows(values_only=True, max_row=1000):  # Limiter à 1000 lignes
                    if any(cell is not None for cell in row):
                        row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                        text += row_text + "\n"
            
            workbook.close()
            return text or "Fichier Excel vide"
        except Exception as e:
            return f"Erreur extraction XLSX: {str(e)}"
    
    @staticmethod
    def _extract_xls_text(file):
        """Extrait le texte d'un fichier Excel .xls"""
        try:
            file.seek(0)
            excel_data = pd.read_excel(file, sheet_name=None, engine='xlrd')
            text = ""
            
            for sheet_name, df in excel_data.items():
                text += f"\n--- Feuille: {sheet_name} ---\n"
                # Limiter le nombre de lignes affichées
                if len(df) > 100:
                    text += df.head(100).to_string(index=False) + f"\n... ({len(df)-100} lignes supplémentaires)\n"
                else:
                    text += df.to_string(index=False) + "\n"
            
            return text or "Fichier Excel vide"
        except Exception as e:
            return f"Erreur extraction XLS: {str(e)}"
    
    @staticmethod
    def _extract_pptx_text(file):
        """Extrait le texte d'un fichier PowerPoint .pptx"""
        try:
            file.seek(0)
            prs = Presentation(file)
            text = ""
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text += f"\n--- Diapositive {slide_num} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
                    
                    # Extraire le texte des tableaux
                    if hasattr(shape, 'has_table') and shape.has_table:
                        text += "\n[Tableau]\n"
                        for row in shape.table.rows:
                            row_text = " | ".join([cell.text.strip() for cell in row.cells])
                            if row_text.strip():
                                text += row_text + "\n"
            
            return text or "Présentation vide"
        except Exception as e:
            return f"Erreur extraction PPTX: {str(e)}"
    
    @staticmethod
    def _extract_csv_text(file):
        """Extrait le contenu d'un fichier CSV"""
        try:
            file.seek(0)
            # Essayer différents encodages
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            
            for encoding in encodings:
                try:
                    file.seek(0)
                    df = pd.read_csv(file, encoding=encoding, nrows=1000)  # Limiter à 1000 lignes
                    result = f"Fichier CSV avec {len(df)} lignes et {len(df.columns)} colonnes:\n\n"
                    result += f"Colonnes: {', '.join(df.columns.tolist())}\n\n"
                    result += df.head(50).to_string(index=False)  # Afficher les 50 premières lignes
                    
                    if len(df) > 50:
                        result += f"\n... ({len(df)-50} lignes supplémentaires)"
                    
                    return result
                except UnicodeDecodeError:
                    continue
            
            return "Impossible de décoder le fichier CSV"
        except Exception as e:
            return f"Erreur extraction CSV: {str(e)}"
    
    @staticmethod
    def _extract_rtf_text(file):
        """Extrait le texte d'un fichier RTF"""
        try:
            file.seek(0)
            content = file.read()
            if isinstance(content, bytes):
                content = content.decode('utf-8', errors='ignore')
            
            # Extraction basique RTF (enlever les balises)
            text = re.sub(r'\\[a-z]+\d*\s?', '', content)
            text = re.sub(r'[{}]', '', text)
            return text.strip() or "Fichier RTF vide"
        except Exception as e:
            return f"Erreur extraction RTF: {str(e)}"
    
    @staticmethod
    def is_image_file(file):
        """Vérifie si le fichier est une image"""
        if not file:
            return False
        
        file_type = getattr(file, 'type', '').lower()
        file_name = getattr(file, 'name', '').lower()
        
        return any(mime_type in file_type for mime_type in IMAGE_MIME_TYPES) or \
               any(file_name.endswith(ext) for ext in IMAGE_EXTENSIONS)
    
    @staticmethod
    def encode_image_to_base64(file):
        """Encode une image en base64 pour l'API Ollama"""
        try:
            file.seek(0)
            if hasattr(file, 'getvalue'):
                return base64.b64encode(file.getvalue()).decode('utf-8')
            else:
                content = file.read()
                return base64.b64encode(content).decode('utf-8')
        except Exception as e:
            raise Exception(f"Erreur encodage base64: {str(e)}")
    
    @staticmethod
    def get_file_hash(file):
        """Génère un hash unique pour éviter les re-traitements"""
        try:
            file.seek(0)
            if hasattr(file, 'getvalue'):
                content = file.getvalue()
            else:
                content = file.read()
            
            return hashlib.md5(content).hexdigest()
        except Exception as e:
            return hashlib.md5(str(e).encode()).hexdigest()
    
    @staticmethod
    def format_file_size(size_bytes):
        """Formate la taille du fichier de manière lisible"""
        if not isinstance(size_bytes, (int, float)) or size_bytes < 0:
            return "Taille inconnue"
        
        if size_bytes < 1024:
            return f"{size_bytes} B"
        elif size_bytes < 1024**2:
            return f"{size_bytes/1024:.1f} KB"
        elif size_bytes < 1024**3:
            return f"{size_bytes/(1024**2):.1f} MB"
        else:
            return f"{size_bytes/(1024**3):.1f} GB"
    
    @staticmethod
    def get_file_type_emoji(file_type, file_name):
        """Retourne l'emoji approprié selon le type de fichier"""
        if not file_name:
            return "📄"
        
        file_name = file_name.lower()
        file_type = file_type.lower() if file_type else ""
        
        if "image" in file_type or any(file_name.endswith(ext) for ext in IMAGE_EXTENSIONS):
            return "🖼️"
        elif "pdf" in file_type or file_name.endswith('.pdf'):
            return "📋"
        elif "word" in file_type or file_name.endswith(('.doc', '.docx')):
            return "📝"
        elif "excel" in file_type or "spreadsheet" in file_type or file_name.endswith(('.xls', '.xlsx')):
            return "📊"
        elif "powerpoint" in file_type or "presentation" in file_type or file_name.endswith(('.ppt', '.pptx')):
            return "📈"
        elif "csv" in file_type or file_name.endswith('.csv'):
            return "📋"
        elif file_name.endswith('.rtf'):
            return "📝"
        elif file_name.endswith('.txt'):
            return "📄"
        else:
            return "📄"